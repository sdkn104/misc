import asyncio
import fnmatch
import io
import os
import smtplib
import xml.etree.ElementTree as ET
from email.header import Header
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from pathlib import Path
from typing import Annotated, Optional

import pypdf
import requests
import yaml
from dotenv import load_dotenv
from pydantic import Field

from agent_framework import MCPStdioTool, tool
from agent_framework.openai import OpenAIChatCompletionClient

# 環境設定
load_dotenv()




# LLMのコンテキスト溢れを防ぐため抽出テキストの上限文字数を設定
MAX_CHARS = 500_000

# ====================================================================================
# === ツール定義 ======================================================================
# ====================================================================================

# --- search_paper_pdf

@tool(approval_mode="never_require")
def search_paper_pdf(
    title: Annotated[str, Field(description="検索する論文タイトル")],
) -> str:
    """論文タイトルからPDF URLを検索する。Semantic Scholar → arXiv の順に試みる。"""
    # 1. Semantic Scholar API (優先)
    try:
        resp = requests.get(
            "https://api.semanticscholar.org/graph/v1/paper/search",
            params={"query": title, "fields": "title,openAccessPdf", "limit": 3},
            timeout=15,
        )
        resp.raise_for_status()
        for paper in resp.json().get("data", []):
            pdf_info = paper.get("openAccessPdf")
            if pdf_info and pdf_info.get("url"):
                return pdf_info["url"]
    except requests.RequestException:
        pass

    # 2. arXiv API (フォールバック)
    try:
        resp = requests.get(
            "https://export.arxiv.org/api/query",
            params={"search_query": f"ti:{title}", "max_results": 3},
            timeout=15,
        )
        resp.raise_for_status()
        ns = {"atom": "http://www.w3.org/2005/Atom"}
        root = ET.fromstring(resp.text)
        for entry in root.findall("atom:entry", ns):
            for link in entry.findall("atom:link", ns):
                if link.get("title") == "pdf":
                    return link.get("href", "")
    except (requests.RequestException, ET.ParseError):
        pass

    return f"PDF URLが見つかりませんでした: '{title}'"

# --- read_pdf_from_url

@tool(approval_mode="never_require")
def read_pdf_from_url(
    url: Annotated[str, Field(description="読み込むPDFファイルのURL")],
    max_pages: Annotated[int, Field(description="最大読み込みページ数。0=全ページ")] = 0,
) -> str:
    """指定URLのPDFをダウンロードしてテキストを抽出する。"""
    try:
        response = requests.get(url, timeout=30)
        response.raise_for_status()
    except requests.Timeout:
        return f"エラー: タイムアウト (30秒) — URL: {url}"
    except requests.HTTPError as e:
        return f"エラー: HTTPエラー {e.response.status_code} — URL: {url}"
    except requests.RequestException as e:
        return f"エラー: ダウンロード失敗 — {e}"

    try:
        reader = pypdf.PdfReader(io.BytesIO(response.content))
    except pypdf.errors.PdfReadError as e:
        return f"エラー: PDFの読み込みに失敗しました — {e}"

    pages = reader.pages
    if max_pages > 0:
        pages = pages[:max_pages]

    parts = []
    for i, page in enumerate(pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            parts.append(f"--- ページ {i} ---\n{text}")

    if not parts:
        return "テキストを抽出できませんでした（画像PDFの可能性があります）。"

    full_text = "\n\n".join(parts)
    clipped = len(full_text) > MAX_CHARS
    if clipped:
        full_text = full_text[:MAX_CHARS] + f"\n\n[... 残り {len(full_text) - MAX_CHARS} 文字は省略されました]"

    return full_text

# --- contents_reader

@tool(approval_mode="never_require")
def contents_reader(
    path: Annotated[str, Field(description=(
        "読み込むファイルのパスまたはURL。"
        "http/https URL、file URL (file:///C:/...)、"
        "Windowsパス (C:\\path\\to\\file)、UNCパス (\\\\server\\share\\file) に対応。"
        "対応形式: PDF, Word (.docx/.doc), Excel (.xlsx/.xls), PowerPoint (.pptx/.ppt), テキスト (.txt/.md)"
    ))],
    engine: Annotated[str, Field(description=(
        "変換エンジン。'fast' (pdfminer/python-docx/openpyxl/python-pptx使用、デフォルト) または "
        "'markitdown' (MarkItDown使用)。.doc/.xls/.ppt の旧形式は markitdown のみ対応。"
    ))] = "fast",
) -> str:
    """PDF・Word・Excel・PowerPoint・テキスト・MarkdownファイルをMarkdownテキストに変換する。"""
    import urllib.parse

    source = path.strip()

    # file: URI → ローカルパスに変換
    if source.lower().startswith("file:"):
        parsed = urllib.parse.urlparse(source)
        local = urllib.parse.unquote(parsed.path)
        if len(local) > 2 and local[0] == "/" and local[2] == ":":
            local = local[1:]
        source = local

    if engine == "markitdown":
        return _read_with_markitdown(source)
    else:
        return _read_with_fast(source)
    

def _read_with_markitdown(source: str) -> str:
    from markitdown import MarkItDown
    try:
        md = MarkItDown(enable_plugins=False)
        if source.lower().startswith(("http://", "https://")):
            result = md.convert_uri(source)
        else:
            result = md.convert(source)
        text = getattr(result, 'text_content', None) or getattr(result, 'markdown', None) or ""
    except Exception as e:
        return f"エラー: 変換に失敗しました — {e}"
    if not text.strip():
        return "テキストを抽出できませんでした。"
    if len(text) > MAX_CHARS:
        text = text[:MAX_CHARS] + f"\n\n[... 残り {len(text) - MAX_CHARS} 文字は省略されました]"
    return text


def _fast_pdf(source) -> str:
    from pdfminer.high_level import extract_text
    return extract_text(source)


def _fast_docx(source) -> str:
    import docx
    doc = docx.Document(source)
    parts = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = para.style.name
        if style.startswith("Heading"):
            tokens = style.split()
            level = tokens[-1] if tokens[-1].isdigit() else "1"
            parts.append(f"{'#' * int(level)} {text}")
        else:
            parts.append(text)
    for table in doc.tables:
        rows = []
        for i, row in enumerate(table.rows):
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
            if i == 0:
                rows.append("|" + " --- |" * len(cells))
        if rows:
            parts.append("\n".join(rows))
    return "\n\n".join(parts)


def _fast_excel(source) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(source, read_only=True, data_only=True)
    parts = []
    for name in wb.sheetnames:
        ws = wb[name]
        rows_data = [r for r in ws.iter_rows(values_only=True) if any(c is not None for c in r)]
        if not rows_data:
            continue
        max_cols = max(len(r) for r in rows_data)
        lines = []
        for i, row in enumerate(rows_data):
            cells = [str(c) if c is not None else "" for c in row]
            cells += [""] * (max_cols - len(cells))
            lines.append("| " + " | ".join(cells) + " |")
            if i == 0:
                lines.append("|" + " --- |" * max_cols)
        parts.append(f"## {name}\n\n" + "\n".join(lines))
    wb.close()
    return "\n\n".join(parts)


def _fast_pptx(source) -> str:
    from pptx import Presentation
    prs = Presentation(source)
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = [f"## スライド {i}"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if len(texts) > 1:
            slides.append("\n\n".join(texts))
    return "\n\n---\n\n".join(slides)


def _read_with_fast(source: str) -> str:
    import urllib.parse

    is_url = source.lower().startswith(("http://", "https://"))

    if is_url:
        ext = Path(urllib.parse.urlparse(source).path).suffix.lower()
        try:
            response = requests.get(source, timeout=30)
            response.raise_for_status()
            content_bytes = response.content
        except requests.Timeout:
            return f"エラー: タイムアウト (30秒) — URL: {source}"
        except requests.HTTPError as e:
            return f"エラー: HTTPエラー {e.response.status_code} — URL: {source}"
        except requests.RequestException as e:
            return f"エラー: ダウンロード失敗 — {e}"
        file_source = io.BytesIO(content_bytes)
    else:
        ext = Path(source).suffix.lower()
        file_source = source
        content_bytes = None

    try:
        if ext == ".pdf":
            text = _fast_pdf(file_source)
        elif ext == ".docx":
            text = _fast_docx(file_source)
        elif ext == ".xlsx":
            text = _fast_excel(file_source)
        elif ext == ".pptx":
            text = _fast_pptx(file_source)
        elif ext in (".txt", ".md", ".markdown", ""):
            if is_url:
                text = content_bytes.decode("utf-8", errors="replace")
            else:
                text = Path(source).read_text(encoding="utf-8", errors="replace")
        elif ext in (".doc", ".xls", ".ppt"):
            return f"エラー: {ext} 形式は fast エンジン非対応です。engine='markitdown' を使用してください。"
        else:
            if is_url:
                text = content_bytes.decode("utf-8", errors="replace")
            else:
                text = Path(source).read_text(encoding="utf-8", errors="replace")
    except Exception as e:
        return f"エラー: 変換に失敗しました — {e}"

    if not text.strip():
        return "テキストを抽出できませんでした。"
    if len(text) > MAX_CHARS:
        text = text[:MAX_CHARS] + f"\n\n[... 残り {len(text) - MAX_CHARS} 文字は省略されました]"
    return text

# --- send_email

_AGENT_CONFIG_PATH = Path(__file__).parent / os.environ.get("AGENT_CONFIG", "agent_config.yaml")

def _load_allowed_patterns() -> list[str]:
    if not _AGENT_CONFIG_PATH.exists():
        return []
    with open(_AGENT_CONFIG_PATH, encoding="utf-8") as f:
        config = yaml.safe_load(f) or {}
    return [p.lower().strip() for p in config.get("allowed_addresses", [])]


def _is_address_allowed(addr: str, patterns: list[str]) -> bool:
    addr_lower = addr.lower().strip()
    return any(fnmatch.fnmatch(addr_lower, p) for p in patterns)


@tool
def send_email(
    to: Annotated[list[str], Field(description="宛先メールアドレスのリスト")],
    subject: Annotated[str, Field(description="メールの件名")],
    body: Annotated[str, Field(description="メールの本文 (プレーンテキスト)")],
    cc: Annotated[Optional[list[str]], Field(description="CCのメールアドレスのリスト")] = None,
    bcc: Annotated[Optional[list[str]], Field(description="BCCのメールアドレスのリスト")] = None,
) -> str:
    """メールを送信する。To/Cc/Bcc は agent_config.yaml で許可されたアドレス/パターンのみ指定可能。"""
    patterns = _load_allowed_patterns()
    if not patterns:
        return f"エラー: 設定ファイル ({_AGENT_CONFIG_PATH}) が見つからないか、allowed_addresses が空です。"

    for field, addrs in [("To", to), ("Cc", cc or []), ("Bcc", bcc or [])]:
        for addr in addrs:
            if not _is_address_allowed(addr, patterns):
                return f"エラー: {field} の '{addr}' は許可されていません。許可パターン: {patterns}"

    smtp_host = os.environ.get("SMTP_HOST")
    if not smtp_host:
        return "エラー: SMTP_HOST が .env に設定されていません。"
    smtp_port = int(os.environ.get("SMTP_PORT", "587"))
    smtp_user = os.environ.get("SMTP_USERNAME", "agent@example.com")
    smtp_pass = os.environ.get("SMTP_PASSWORD", "")
    smtp_from = os.environ.get("SMTP_FROM", smtp_user)
    use_starttls = os.environ.get("SMTP_USE_STARTTLS", "true").lower() in ("true", "1", "yes")

    msg = MIMEMultipart()
    msg["From"] = smtp_from
    msg["To"] = ", ".join(to)
    msg["Subject"] = Header(subject, "utf-8")
    if cc:
        msg["Cc"] = ", ".join(cc)
    msg.attach(MIMEText(body, "plain", "utf-8"))

    all_recipients = to + (cc or []) + (bcc or [])
    try:
        if use_starttls:
            server = smtplib.SMTP(smtp_host, smtp_port, timeout=30)
            server.starttls()
        else:
            server = smtplib.SMTP_SSL(smtp_host, smtp_port, timeout=30)
        if smtp_user and smtp_pass:
            server.login(smtp_user, smtp_pass)
        server.sendmail(smtp_from, all_recipients, msg.as_string())
        server.quit()
    except smtplib.SMTPException as e:
        return f"エラー: SMTP エラー — {e}"
    except OSError as e:
        return f"エラー: 接続失敗 — {e}"

    summary = f"To: {', '.join(to)}"
    if cc:
        summary += f" / Cc: {', '.join(cc)}"
    if bcc:
        summary += f" / Bcc: {', '.join(bcc)}"
    return f"送信完了。件名: 「{subject}」、{summary}"


# --- search_active_directory

@tool(approval_mode="never_require")
def search_active_directory(
    query: Annotated[str, Field(description=(
        "検索クエリ。名前・メール・アカウント名などのキーワード、または LDAP フィルタ文字列 "
        "(例: '(sAMAccountName=jdoe)') を指定する。"
    ))],
    search_base: Annotated[Optional[str], Field(description=(
        "LDAP 検索ベース DN (例: 'ou=Users,dc=example,dc=com')。"
        "省略時は環境変数 LDAP_SEARCH_BASE を使用。"
    ))] = None,
    limit: Annotated[int, Field(description="返す最大件数")] = 50,
) -> str:
    """Active Directory を LDAP で検索し、ユーザ・グループ情報を返す。"""
    try:
        import ldap3
        import ldap3.utils.conv
    except ImportError:
        return "エラー: ldap3 パッケージがインストールされていません。pip install ldap3 を実行してください。"

    server_url = os.environ.get("LDAP_SERVER")
    if not server_url:
        return "エラー: LDAP_SERVER が .env に設定されていません。"

    port = int(os.environ.get("LDAP_PORT", "389"))
    use_ssl = os.environ.get("LDAP_USE_SSL", "false").lower() in ("true", "1", "yes")
    bind_dn = os.environ.get("LDAP_BIND_DN", "")
    bind_pw = os.environ.get("LDAP_BIND_PASSWORD", "")
    base = search_base or os.environ.get("LDAP_SEARCH_BASE", "")
    if not base:
        return "エラー: LDAP_SEARCH_BASE が .env に設定されていません。"

    if query.startswith("("):
        ldap_filter = query
    else:
        esc = ldap3.utils.conv.escape_filter_chars(query)
        ldap_filter = (
            f"(|(displayName=*{esc}*)"
            f"(mail=*{esc}*)"
            f"(sAMAccountName=*{esc}*)"
            f"(cn=*{esc}*))"
        )

    attrs = [
        "displayName", "sAMAccountName", "mail", "department",
        "title", "telephoneNumber", "memberOf", "objectClass",
        "cn", "description",
    ]

    try:
        server = ldap3.Server(server_url, port=port, use_ssl=use_ssl, get_info=ldap3.NONE, connect_timeout=10)
        conn = ldap3.Connection(server, user=bind_dn, password=bind_pw, auto_bind=True, receive_timeout=30)
    except ldap3.core.exceptions.LDAPException as e:
        return f"エラー: LDAP 接続失敗 — {e}"

    try:
        conn.search(
            search_base=base,
            search_filter=ldap_filter,
            search_scope=ldap3.SUBTREE,
            attributes=attrs,
            size_limit=limit,
        )
    except ldap3.core.exceptions.LDAPException as e:
        conn.unbind()
        return f"エラー: LDAP 検索失敗 — {e}"

    results = []
    for entry in conn.entries:
        obj: dict = {"dn": entry.entry_dn}
        for attr in attrs:
            try:
                val = entry[attr].value
                if val is not None:
                    obj[attr] = val
            except Exception:
                pass
        results.append(obj)
    conn.unbind()

    if not results:
        return f"検索結果が見つかりませんでした。フィルタ: {ldap_filter}"

    import json
    return json.dumps(results, ensure_ascii=False, default=str, indent=2)


# ====================================================================================
# === MCPサーバーを定義　===============================================================
# ====================================================================================

fs_root = Path(os.environ.get("AGENT_FS_ROOT", Path.cwd().resolve()))
filesystem_mcp = MCPStdioTool(
    name="filesystem",
    command="cmd",
    args=["/c", os.getcwd()+"/node_modules/.bin/mcp-server-filesystem.cmd", str(fs_root), str(Path.home() / "Downloads")],
    #command="D:\NoSync\misc\Comp\ChatUI/node_modules/.bin/mcp-server-filesystem.cmd",
    #command="npx",
    #args=["-y", "@modelcontextprotocol/server-filesystem", str(fs_root)],
)
playwright_mcp = MCPStdioTool(
    name="playwright",
    command="cmd",
    args=["/c", os.getcwd()+"/node_modules/.bin/playwright-mcp.cmd", "--browser=msedge", "--allowed-origins=https://*"],
    #command="npx",
    #args=["-y", "@playwright/mcp", "--browser=msedge", "--allowed-origins=https://*,http://*"],
    #args=["-y", "@playwright/mcp", "--browser=msedge", "--allowed-origins=https://*"],
    #args=["-y", "@playwright/mcp", "--browser=msedge"],
)
# markitdown_mcp = MCPStdioTool(
#     name="markitdown",
#     #command=os.getcwd()+r"\myenv\Scripts\markitdown-mcp.exe",
#     #args=[],
#     command="node.exe",
#     args=[os.getcwd()+r"\node_modules\markitdown-mcp-npx\bin\markitdown-mcp-npx.js"],
#     #command="cmd",
#     #args=["/c", os.getcwd()+"\node_modules\.bin\markitdown-mcp-npx.cmd"],
# )


class MCPServerTool:
    """MCPサーバーを起動/シャットダウンするツールのクラス。"""
    def __init__(self, mcp_tools):
        self.mcp_tools = mcp_tools

    def tools(self):
        """MCPサーバーが提供するツールのリストを返す。"""
        return self.mcp_tools

    async def start(self):
        """MCPサーバーを起動する。"""
        for tool in self.mcp_tools:
            await tool.connect()
            print(f"{tool.name} MCP server started")

    async def shutdown(self):
        """MCPサーバーをシャットダウンする。"""
        for tool in self.mcp_tools:
            await tool.close()

mcp_servers = MCPServerTool([filesystem_mcp, playwright_mcp])

# ====================================================================================
# === エージェント ====================================================================
# ====================================================================================

def _display_content(content) -> None:
    """蓄積済み content を type ごとにまとめて表示する。"""
    def _trunc(s: str, n: int) -> str:
        s = str(s)
        return s[:n] + "..." if len(s) > n else s

    match content.type:
        case "text":
            if content.text:
                print(f"\n{("-" * 60)}\n{content.text}", flush=True)
        case "text_reasoning":
            if content.text:
                print(f"[Thinking] {content.text}", flush=True)
        case "function_call":
            print(f"\n[Tool] {content.name}({_trunc(content.arguments, 120)})", flush=True)
        case "function_result":
            print(f"  → {_trunc(content.result, 80)}", flush=True)
        case "mcp_server_tool_call":
            print(f"\n[Browser] {content.name}({_trunc(content.arguments, 120)})", flush=True)
        case "mcp_server_tool_result":
            print(f"  → {_trunc(content.result, 80)}", flush=True)


async def agent_run(agent, user_input, session, model: str | None = None, effort: str | None = None):
    """agent.run をラップし、同一 type のチャンクを蓄積して yield する async generator。"""
    options: dict = {}
    if model:
        options["model"] = model
    if effort:
        options["reasoning"] = {"effort": effort}
    acc = None
    async for update in agent.run(user_input, stream=True, session=session, options=options or None):
        for content in update.contents:
            if acc is None:
                acc = content
            elif acc.type == content.type:
                try:
                    acc = acc + content
                except Exception:
                    yield acc
                    acc = content
            else:
                yield acc
                acc = content
    if acc is not None:
        yield acc


def _load_instructions(filename: str = "SYSTEM.md", user: str = "unknown") -> str:
    """スクリプトと同じディレクトリの指定ファイルを読み込んで返す。"""
    path = Path(__file__).parent / filename
    return path.read_text(encoding="utf-8") + f"\n\n## ユーザ情報\n- 情報システムID： {user}\n"


def build_client(
    model: str | None = None,
    endpoint: str | None = None,
    api_key: str | None = None,
    api_version: str | None = None,
) -> OpenAIChatCompletionClient:
    endpoint = endpoint or os.environ.get("AZURE_OPENAI_ENDPOINT")
    model = model or os.environ.get("AZURE_OPENAI_DEPLOYMENT_NAME")
    api_version = api_version or os.environ.get("OPENAI_API_VERSION")
    api_key = api_key or os.environ.get("AZURE_OPENAI_API_KEY")

    if not endpoint or not model:
        raise RuntimeError(
            "AZURE_OPENAI_ENDPOINT と AZURE_OPENAI_DEPLOYMENT_NAME を .env または環境変数に設定してください。"
        )

    return OpenAIChatCompletionClient(
        model=model,
        azure_endpoint=endpoint,
        api_version=api_version,
        api_key=api_key,
    )


# ====================================================================================
# メイン（単独実行時）
# ====================================================================================

async def main() -> None:

    # === MCP servers ===============================================================
    await mcp_servers.start()

    # === build agent and run =======================================================
    import getpass
    user = getpass.getuser()
    client = build_client()
    agent = client.as_agent(
        name="DocAgent",
        instructions=_load_instructions(user=user),
        tools=[search_paper_pdf, read_pdf_from_url, contents_reader, send_email, search_active_directory] + mcp_servers.tools(),
    )

    print("Doc Agent 起動中。ファイル名やURLを含む質問を入力してください。")
    print("例1: Attention Is All You Need を要約して")
    print("例2: https://arxiv.org/abs/1706.03762 をブラウザで開いて要約して")
    print("終了するには 'quit' または 'exit' を入力してください。\n")

    session = agent.create_session()
    while True:
        try:
            user_input = input("You: ").strip()
        except (EOFError, KeyboardInterrupt):
            break

        if not user_input:
            continue
        if user_input.lower() in ("quit", "exit"):
            break

        #print("\nAgent:", flush=True)
        async for content in agent_run(agent, user_input, session):
            _display_content(content)
        print("\n")

    await mcp_servers.shutdown()

if __name__ == "__main__":
    asyncio.run(main())
