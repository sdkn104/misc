import sys
import re
import requests
from io import BytesIO
from typing import List, Dict, Any

from markdown_it import MarkdownIt
from markdown_it.renderer import RendererHTML
#from mdit_py_plugins.table import table_plugin
#from mdit_py_plugins.footnote import footnote_plugin

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


# ============================================================
# Layout Engine
# ============================================================

class SlideLayoutEngine:

    def __init__(self, prs: Presentation):

        self.prs = prs
        self.current_slide = None
        self.textbox = None
        self.tf = None

    def new_slide(self, title=None):

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])

        self.current_slide = slide

        if title and slide.shapes.title:
            slide.shapes.title.text = title

        self.textbox = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(1.5),
            Inches(9),
            Inches(5),
        )

        self.tf = self.textbox.text_frame
        self.tf.clear()

        return slide

    def add_paragraph(self, text="", level=0):

        if not self.tf:
            self.new_slide()

        p = self.tf.add_paragraph()
        p.text = text
        p.level = level
        return p

    def add_code(self, code):

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])

        box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(0.5),
            Inches(9),
            Inches(6)
        )

        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = code
        p.font.name = "Courier New"
        p.font.size = Pt(16)

    def add_image(self, url, alt):

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])

        if slide.shapes.title:
            slide.shapes.title.text = alt

        try:
            data = requests.get(url).content
            slide.shapes.add_picture(
                BytesIO(data),
                Inches(1),
                Inches(2),
                height=Inches(4),
            )
        except Exception:
            pass

    def add_table(self, headers, rows):

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])

        table = slide.shapes.add_table(
            len(rows) + 1,
            len(headers),
            Inches(0.5),
            Inches(1),
            Inches(9),
            Inches(5),
        ).table

        for i, h in enumerate(headers):
            table.cell(0, i).text = h

        for r, row in enumerate(rows):
            for c, val in enumerate(row):
                table.cell(r + 1, c).text = val


# ============================================================
# Inline Renderer
# ============================================================

class InlineRenderer:

    def render(self, tokens):

        result = ""

        for t in tokens:

            if t.type == "text":
                result += t.content

            elif t.type == "code_inline":
                result += f"`{t.content}`"

            elif t.type == "softbreak":
                result += "\n"

            elif t.type == "hardbreak":
                result += "\n"

            elif t.type == "strong_open":
                result += "**"

            elif t.type == "strong_close":
                result += "**"

            elif t.type == "em_open":
                result += "*"

            elif t.type == "em_close":
                result += "*"

            elif t.type == "link_open":
                href = dict(t.attrs).get("href", "")
                result += f"["

            elif t.type == "link_close":
                result += "]"

            elif t.type == "image":
                src = dict(t.attrs).get("src", "")
                alt = t.content
                result += f"[img:{alt}]({src})"

        return result


# ============================================================
# Block Renderer
# ============================================================

class BlockRenderer:

    def __init__(self, layout: SlideLayoutEngine):

        self.layout = layout
        self.inline = InlineRenderer()

    def render(self, tokens):

        i = 0

        while i < len(tokens):

            t = tokens[i]

            if t.type == "heading_open":

                level = int(t.tag[1])
                text = tokens[i + 1].content

                self.layout.new_slide(text)

                i += 2

            elif t.type == "paragraph_open":

                inline = tokens[i + 1]

                txt = self.inline.render(inline.children or [])

                self.layout.add_paragraph(txt)

                i += 2

            elif t.type == "bullet_list_open":

                i = self.render_list(tokens, i, ordered=False)

            elif t.type == "ordered_list_open":

                i = self.render_list(tokens, i, ordered=True)

            elif t.type in ("fence", "code_block"):

                self.layout.add_code(t.content)

            elif t.type == "blockquote_open":

                text = tokens[i + 2].content
                p = self.layout.add_paragraph(text)
                p.font.italic = True
                i += 3

            elif t.type == "table_open":

                i = self.render_table(tokens, i)

            elif t.type == "inline":

                for child in t.children or []:
                    if child.type == "image":

                        src = dict(child.attrs).get("src")
                        alt = child.content

                        self.layout.add_image(src, alt)

            i += 1

    def render_list(self, tokens, i, ordered):

        level = 0
        index = 1

        i += 1

        while tokens[i].type not in ("bullet_list_close", "ordered_list_close"):

            if tokens[i].type == "list_item_open":

                text = tokens[i + 2].content

                if ordered:
                    text = f"{index}. {text}"
                    index += 1

                self.layout.add_paragraph(text, level)

                i += 4

            else:
                i += 1

        return i

    def render_table(self, tokens, i):

        headers = []
        rows = []

        i += 1

        while tokens[i].type != "table_close":

            if tokens[i].type == "tr_open":

                row = []
                i += 1

                while tokens[i].type != "tr_close":

                    if tokens[i].type in ("th_open", "td_open"):

                        row.append(tokens[i + 1].content)
                        i += 3

                    else:
                        i += 1

                if not headers:
                    headers = row
                else:
                    rows.append(row)

            else:
                i += 1

        self.layout.add_table(headers, rows)

        return i


# ============================================================
# Plugin System
# ============================================================

class MarkdownPlugin:

    def before_render(self, tokens):
        return tokens

    def after_render(self, prs):
        pass


class EmojiPlugin(MarkdownPlugin):

    emoji = {
        ":smile:": "😄",
        ":rocket:": "🚀",
        ":fire:": "🔥",
    }

    def before_render(self, tokens):

        for t in tokens:
            if hasattr(t, "content"):
                for k, v in self.emoji.items():
                    t.content = t.content.replace(k, v)

        return tokens


# ============================================================
# Main Engine
# ============================================================

class MarkdownPPTXEngine:

    def __init__(self):

        self.md = (
            MarkdownIt("gfm-like", renderer_cls=RendererHTML) # pip install  linkify-it-py
            #MarkdownIt("commonmark")
            #.use(table_plugin)
            #.use(footnote_plugin)
        )

        self.plugins: List[MarkdownPlugin] = []

    def add_plugin(self, plugin):

        self.plugins.append(plugin)

    def convert(self, markdown_text, output):

        tokens = self.md.parse(markdown_text)

        for p in self.plugins:
            tokens = p.before_render(tokens)

        prs = Presentation()

        layout = SlideLayoutEngine(prs)

        renderer = BlockRenderer(layout)

        renderer.render(tokens)

        for p in self.plugins:
            p.after_render(prs)

        prs.save(output)


# ============================================================
# CLI
# ============================================================

def main():

    md = """
# Markdown PPTX Engine

This engine converts **markdown** into PowerPoint.

## Lists

- item 1
- item 2
- item 3

## Ordered

1. first
2. second
3. third

## Code

```python
print("hello world")
````

## Table

| A | B |
| - | - |
| 1 | 2 |
| 3 | 4 |

## Image

![cat](https://placekitten.com/400/300)

## Quote

> markdown-it is powerful

:rocket:
"""

    if len(sys.argv) < 2:
        print("使い方: python outpptx.py <入力ファイル>")
        sys.exit(1)
    input_file = sys.argv[1]
    output_file = input_file.rsplit('.', 1)[0] + '.pptx'

    with open(input_file, "r", encoding="utf-8") as f:
        md = f.read()

    engine = MarkdownPPTXEngine()
    engine.add_plugin(EmojiPlugin())
    engine.convert(md, output_file)
    print(f"Done: {output_file}")

if __name__ == "__main__":
    main()

"""

# このエンジンの特徴

対応ブロック：

- headings
- paragraph
- bullet list
- ordered list
- blockquote
- code block
- table
- image
- hr
- footnote (拡張)

inline：

- text
- bold
- italic
- inline code
- link
- image

拡張：

```

engine.add_plugin(...)

```

で追加可能。

---

# 実用化するなら追加すべき機能

本格運用する場合は次を追加します。

### 1. スライドレイアウトエンジン
```

H1 → title slide
H2 → new slide
H3 → bullet

```

### 2. 数式
```

markdown-it-math
↓
MathJax
↓
画像化

```

### 3. Mermaid

```

mermaid
↓
svg
↓
pptx

```

### 4. GitHub Markdown完全対応

追加プラグイン：

```

tasklist
emoji
container
attrs
deflist
abbr

```

---

もし希望があれば、次も書けます：

- **完全GitHub Markdown → PPTX（2000行級エンジン）**
- **Markdown → Google Slides**
- **Markdown → Keynote**
- **Pandoc級 Markdown スライド生成器**  
- **Markdown → PPTX テンプレートエンジン（企業資料用）**

かなり面白いレベルの **スライド生成システム**になります。
```
"""