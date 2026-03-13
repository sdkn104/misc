from pptx import Presentation
from markdown_it import MarkdownIt
from pprint import pprint
from markdown_it.renderer import RendererHTML
import sys
import os

#md = MarkdownIt()
#md = MarkdownIt("commonmark").enable('table')
md = MarkdownIt("gfm-like", renderer_cls=RendererHTML) # pip install  linkify-it-py
renderer = RendererHTML()

def parse_marp(markdown_text):
    slides = markdown_text.split("\n---\n")
    parsed = []

    for slide in slides:
        tokens = md.parse(slide)
        for t in tokens:
            print(t.as_dict())

        print(md.render(slide))
        #result = renderer.render(tokens, {}, "")
        #print(result)


        title = None
        bullets = []
        tables = []

        i = 0
        while i < len(tokens):
            t = tokens[i]

            # title
            if t.type == "heading_open":
                title = tokens[i+1].content

            # bullet list
            if t.type == "inline" and tokens[i-1].type == "list_item_open":
                bullets.append(t.content)

            # table
            if t.type == "table_open":
                table = []
                i += 1
                row = []
                while tokens[i].type != "table_close":
                    if tokens[i].type == "inline":
                        row.append(tokens[i].content)
                    if tokens[i].type == "tr_close":
                        table.append(row)
                        row = []
                    i += 1
                tables.append(table)

            i += 1

        parsed.append({
            "title": title,
            "bullets": bullets,
            "tables": tables
        })

    return parsed


def marp_to_pptx(markdown_text, output="slides.pptx"):
    prs = Presentation()

    slides = parse_marp(markdown_text)
    pprint(slides)

    for s in slides:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)

        title = slide.shapes.title
        body = slide.placeholders[1]

        if s["title"]:
            title.text = s["title"]

        tf = body.text_frame
        tf.clear()

        # bullets
        for b in s["bullets"]:
            p = tf.add_paragraph()
            p.text = b
            p.level = 0

        # tables
        for table in s["tables"]:
            rows = len(table)
            cols = len(table[0])

            left = 0
            top = 3000000
            width = 9000000
            height = 3000000

            t = slide.shapes.add_table(rows, cols, left, top, width, height).table

            for r in range(rows):
                for c in range(cols):
                    t.cell(r,c).text = table[r][c]

    print(prs)
    prs.save(output)


# ===== example =====

marp_md = """
# Sales Report

- Revenue increased
- Marketing improved

| region | sales |
|------|------|
| US | 100 |
| EU | 80 |

---

# Next Plan

- Launch campaign
- Expand team
"""



def main():
    # pandocのパスをPATHに追加
    #pandoc_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'pandoc-3.9-windows-x86_64', 'pandoc-3.9')
    #os.environ["PATH"] = pandoc_dir + os.pathsep + os.environ.get("PATH", "")

    if len(sys.argv) < 2:
        print("使い方: python outpptx.py <入力ファイル>")
        sys.exit(1)
    input_file = sys.argv[1]
    output_file = input_file.rsplit('.', 1)[0] + '.pptx'

    with open(input_file, "r", encoding="utf-8") as f:
        marp_md = f.read()
    marp_to_pptx(marp_md, output_file)
    print(f"Done: {output_file}")

if __name__ == "__main__":
    main()

